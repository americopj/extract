from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import fitz
from io import BytesIO
import uvicorn

app = FastAPI()


@app.get("/")
def health():
    return {"status": "ok"}


def extract_pptx(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def extract_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    return "\n\n".join(paragraphs)


def extract_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_text = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            pages_text.append(f"[Página {i}]\n{text}")
    doc.close()
    return "\n\n".join(pages_text)


def extract_xlsx(file_bytes):
    wb = load_workbook(BytesIO(file_bytes), read_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            if any(row_text):
                rows.append(" | ".join(row_text))
        if rows:
            sheets_text.append(f"[Planilha: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets_text)


@app.post("/extract")
async def extract_file(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        filename = file.filename.lower()
        
        if filename.endswith(".pptx"):
            text = extract_pptx(contents)
            file_type = "pptx"
        elif filename.endswith(".docx"):
            text = extract_docx(contents)
            file_type = "docx"
        elif filename.endswith(".pdf"):
            text = extract_pdf(contents)
            file_type = "pdf"
        elif filename.endswith(".xlsx"):
            text = extract_xlsx(contents)
            file_type = "xlsx"
        else:
            return JSONResponse({
                "success": False,
                "error": f"Formato não suportado: {filename}"
            }, status_code=400)
        
        return JSONResponse({
            "success": True,
            "text": text,
            "file_type": file_type,
            "filename": file.filename
        })
    
    except Exception as e:
        return JSONResponse({
            "success": False,
            "error": str(e)
        }, status_code=400)


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
